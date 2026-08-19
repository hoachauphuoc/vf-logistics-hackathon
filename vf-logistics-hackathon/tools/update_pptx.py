"""Update VF_Logistics_Presentation.pptx for the Refinement Phase.

Two kinds of change:

1. Corrections. The deck was generated 2026-08-02 and states things that are no
   longer true, and in two places states things that were never true:
     * "10,001 shipments / $53,042,546" -> the data holds 10,017 / $52,869,484
     * "Four-field validation -> 0/25/50/75/100" -> validation is six rules and
       the score is ROUND(100 * (6 - failed) / 6)
     * AI_DECISION -> the column is actually AI_RECOMMENDED_ACTION. A judge who
       copies the slide's field name into a query gets "invalid identifier".
     * Mendix is described as the live review interface, but the Refinement Phase
       removed it in favour of Streamlit-in-Snowflake.
     * "Cortex Agent" is listed as a platform feature, but SHOW AGENTS returns no
       object, so the claim is not backed by anything a judge can inspect.

2. Three new slides covering the Refinement Phase itself, which the deck did not
   mention at all even though it is the entire deliverable of this round.

Every figure comes from tools/deck_figures.py, computed off the 2026-08-19 backup
that the judging account will be restored from -- not from the older deck and not
from the pre-refinement state still on AYUGBCE-JX50275.

New slides are produced by deep-copying an existing slide's XML rather than by
calling add_slide() on a blank layout. The deck was exported from Google Slides and
every slide carries its own background picture; a fresh blank slide would render
with no background and obviously different type styling.
"""

import copy
import sys
from pathlib import Path

from pptx import Presentation

PATH = Path(r"C:\Users\phuochoa\Mendix\VF_Logistics_Portal-main_2"
            r"\snowflake-backend\VF_Logistics_Presentation.pptx")

# ---------------------------------------------------------------------------
# Verified figures (tools/deck_figures.py against backup_2026-08-19/data)
# ---------------------------------------------------------------------------
BILLS = "10,017"
CHARGES = "$52,869,484"
CHARGES_SHORT = "$52.9M"

# ---------------------------------------------------------------------------
# Plain text replacements, applied to every run in the deck.
# Longest-first so a shorter key cannot pre-empt a longer overlapping one.
# ---------------------------------------------------------------------------
REPLACEMENTS = [
    ("10,001 shipments / $53,042,546 cargo charges",
     f"{BILLS} shipments / {CHARGES} cargo charges"),
    ("10,001 shipments / $53.0M+", f"{BILLS} shipments / {CHARGES_SHORT}"),
    ("10,001", BILLS),

    # Confidence scoring: six rules, not four fields.
    ("Four-field validation -> score of 0 / 25 / 50 / 75 / 100",
     "Six deterministic rules -> ROUND(100 x (6 - failed) / 6)"),

    # Mendix removed as an interface (Refinement Phase, Fix 2).
    ("PUT file://*.pdf @LOGISTICS_STAGE - many files in one command "
     "(the Mendix UI uploads one file at a time)",
     "PUT file://*.pdf @LOGISTICS_STAGE - many files in one command; "
     "Snowsight stage upload for ad-hoc drops"),
    ("7. Human review (Mendix UI)", "7. Human review (Streamlit)"),
    ("Reviewer opens the record in Mendix and corrects fields on the "
     "AI-extracted form",
     "Reviewer opens the record in the Streamlit Review & Edit panel and "
     "corrects fields on the AI-extracted form"),
    ("Mendix writes the UPDATE back to BILL_OF_LADING_EXTRACTED "
     "(verified live in QUERY_HISTORY) and sets STATUS = Synced_To_SAP",
     "Streamlit writes the UPDATE back to BILL_OF_LADING_EXTRACTED via "
     "REVIEW_DOCUMENT and sets STATUS = Synced_To_SAP"),
    ("One entry point: WORKFLOW_INGEST_AND_DECIDE (CLI, Mendix chat, Snowpark, "
     "or the stream-driven task)",
     "One entry point: WORKFLOW_INGEST_AND_DECIDE (CLI, Streamlit, Snowpark, "
     "or the stream-driven task)"),
    ("Key Snowflake features: Cortex Agent, Cortex AI (COMPLETE), Marketplace, "
     "Streams/Tasks, Streamlit, Snowpark.",
     "Key Snowflake features: Cortex AI (COMPLETE + PARSE_DOCUMENT), Cortex "
     "Search, Marketplace, Streams/Tasks, Dynamic Tables, Streamlit, Snowpark."),
    ("CLI, Cortex Agent (NL), Python/Snowpark, Streamlit",
     "CLI, Streamlit (EN/JA/VN), Python/Snowpark, SQL"),
    ("Monitoring UI", "Monitoring UI"),
    ("Native Streamlit dashboard with verified KPIs",
     "Native Streamlit dashboard, trilingual EN/JA/VN, with verified KPIs"),
    ("2) Uses Python + Java: Snowpark/Streamlit (Python) + Mendix integration "
     "(Java)",
     "2) Python throughout: Snowpark procedures, Streamlit-in-Snowflake UI, "
     "and the pre-deploy test tooling"),
    ("3) Uses Snowflake platform: Cortex AI, Cortex Agent, Marketplace, "
     "Streams/Tasks",
     "3) Uses Snowflake platform: Cortex AI, Cortex Search, Marketplace, "
     "Streams/Tasks, Dynamic Tables"),
    ("Mendix MVP: vflogisticsportal-sandbox.mxapps.io",
     "Streamlit is the only interface (Mendix retired in refinement)"),
    ("Cost control: tasks and Cortex Search suspended by default, warehouse "
     "auto-suspend at 60s",
     "Cost control: 7 tasks are stream-gated so idle runs are SKIPPED, Cortex "
     "Search suspended, warehouse auto-suspend at 60s"),
    ("Deck generated from PRESENTATION_OUTLINE.md - updated 2026-08-02",
     "Deck generated from PRESENTATION_OUTLINE.md - updated 2026-08-19 "
     "(Refinement Phase)"),
]

# Whole-string replacements. Applied only when a paragraph or cell equals the key
# exactly. AI_DECISION has to be handled this way: it is a substring of
# AI_DECISION_REASON, which is already the correct column name, so a substring
# replace would corrupt it into AI_RECOMMENDED_ACTION_REASON.
EXACT_REPLACEMENTS = {
    "AI_DECISION": "AI_RECOMMENDED_ACTION",
    "Action chosen by the model": "Action chosen by the model",
}

# ---------------------------------------------------------------------------
# New slides. Each is (source_slide_index_to_clone, title, table_or_bullets)
# Cloned from slide 14 (title + 3-column table) or slide 13 (title + bullets),
# both 0-based below.
# ---------------------------------------------------------------------------
TABLE_TEMPLATE = 13   # slide 14: "Technical Execution Highlights" + 6x3 table
BULLET_TEMPLATE = 12  # slide 13: "Technical Execution Highlights" + bullets

NEW_TABLE_SLIDES = [
    ("Refinement Phase - Evaluator Feedback Addressed", [
        ["Feedback", "What was wrong", "What it is now"],
        ["1. File ingestion & admin errors",
         "SELECT INTO raised 'expects exactly 1 row' - 20 duplicate ALERT_ID "
         "from autoincrement colliding with seeded ids",
         "Root-caused and fixed; ids are sequence-backed. "
         "WORKFLOW_INGEST_AND_DECIDE returns COMPLETED"],
        ["2. Merge sandbox portal into Streamlit",
         "Operator work lived in an external Mendix portal",
         "Mendix retired. Process, review, approve/reject and SAP posting are "
         "native Streamlit, plus a shared ui.py theme"],
        ["3. Chat with session history",
         "Chat was in-memory; a reload lost the conversation",
         "Persisted in CHAT_SESSION / CHAT_MESSAGE via 6 owner-rights "
         "procedures; survives reload and is scoped per CURRENT_USER()"],
    ]),
    ("Refinement Phase - Defects Found Beyond the Feedback", [
        ["Defect", "Why it mattered", "Fix"],
        ["Compliance reported a false pass",
         "BATCH_CHECK_COMPLIANCE evaluated no rules and returned "
         "{passed:10017, failed:0} for shipments never examined",
         f"Rules now evaluated set-based. Real result: 8,666 pass / "
         f"1,351 fail (13.5% of {BILLS})"],
        ["Fraud detector had latched off",
         "Backpressure counted every open alert ever raised, so once the "
         "backlog crossed the limit it never scanned again",
         "Queue is a 7-day rolling window; HIGH-severity rules bypass the "
         "gate entirely so a MEDIUM backlog cannot hide a HIGH fraud"],
        ["Cortex spend under-reported",
         "Extraction made 2 Cortex calls per document and logged neither, so "
         "AI FinOps showed a fraction of real usage",
         "Both calls logged with measured latency and real token counts "
         "from COUNT_TOKENS"],
        ["Unenforced primary keys",
         "Snowflake does not enforce PK, and 'autoincrement noorder' hands out "
         "ids from non-monotonic cached ranges",
         "12 sequences replace autoincrement on every table loaded with "
         "explicit ids; 0 duplicates remain"],
    ]),
]

NEW_BULLET_SLIDES = [
    ("Refinement Phase - How It Was Verified", [
        "End-to-end test: 10 purpose-built PDFs covering clean, missing-field, "
        "placeholder, implausible-weight, stale-date and carrier-mismatch cases "
        "- 10/10 pass, 50/50 field extractions exact",
        "System restored to its pre-test state afterwards and proven so: "
        "33/33 row counts and 32/33 content hashes identical (the one "
        "difference is a dynamic table containing CURRENT_TIMESTAMP by design)",
        "Trilingual UI gated by tooling, not by eye: 348 keys x 3 languages, "
        "zero inline language conditionals, zero unresolved keys",
        "Pre-deploy gate runs pyflakes and loads all 7 pages against a stubbed "
        "Streamlit - built after a NameError reached the deployed app",
        "Every fix validated by a test written to fail: where the obvious test "
        "passed for the wrong reason, the defect was reintroduced to confirm "
        "the test could catch it",
        f"AI decisions are real and differentiated: CLEAR 270 / ESCALATE 55 / "
        f"BLOCK 43 across 368 assessed alerts",
    ]),
]

INSERT_AFTER = 13   # 0-based: after slide 14 (Technical Execution Highlights)


def replace_text_everywhere(prs):
    """Apply REPLACEMENTS to every run, and to table cells, in the deck."""
    hits = {}

    def fix_text_frame(tf):
        for para in tf.paragraphs:
            if not para.runs:
                continue
            full = "".join(r.text for r in para.runs)
            new = full

            exact = EXACT_REPLACEMENTS.get(full.strip())
            if exact is not None and exact != full.strip():
                new = exact
                hits[full.strip()] = hits.get(full.strip(), 0) + 1
            else:
                for old, rep in sorted(REPLACEMENTS, key=lambda x: -len(x[0])):
                    if old in new and old != rep:
                        new = new.replace(old, rep)
                        hits[old] = hits.get(old, 0) + 1
            if new != full:
                # Collapse into the first run so mid-word run splits (very
                # common in Google Slides exports) cannot corrupt the result.
                para.runs[0].text = new
                for extra in para.runs[1:]:
                    extra.text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                fix_text_frame(shape.text_frame)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        fix_text_frame(cell.text_frame)
    return hits


def clone_slide(prs, index):
    """Deep-copy slide `index` and return the new slide (appended at the end)."""
    src = prs.slides[index]
    dst = prs.slides.add_slide(src.slide_layout)

    # Drop placeholders the layout added; the cloned shapes carry everything.
    for shp in list(dst.shapes):
        shp._element.getparent().remove(shp._element)

    for shp in src.shapes:
        dst.shapes._spTree.append(copy.deepcopy(shp._element))
    return dst


def set_table(shape, data):
    """Resize the cloned table to len(data) rows and write the cell text."""
    tbl = shape.table
    want, have = len(data), len(tbl.rows)

    # python-pptx table rows do not accept negative indexing.
    while len(tbl.rows) > want:
        tbl._tbl.remove(tbl.rows[len(tbl.rows) - 1]._tr)
    while len(tbl.rows) < want:
        tbl._tbl.append(copy.deepcopy(tbl.rows[len(tbl.rows) - 1]._tr))

    for r_i, row_vals in enumerate(data):
        for c_i, val in enumerate(row_vals):
            if c_i >= len(tbl.columns):
                continue
            cell = tbl.cell(r_i, c_i)
            para = cell.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = val
                for extra in para.runs[1:]:
                    extra.text = ""
            else:
                para.add_run().text = val


def set_title(slide, text):
    """Write `text` into the slide's first non-empty text box (the title)."""
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            para = shape.text_frame.paragraphs[0]
            if para.runs:
                para.runs[0].text = text
                for extra in para.runs[1:]:
                    extra.text = ""
                return True
    return False


def set_bullets(slide, lines):
    """Fill the largest multi-paragraph text box with `lines`."""
    best, best_n = None, -1
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        n = sum(1 for p in shape.text_frame.paragraphs if p.runs)
        if n > best_n:
            best, best_n = shape, n
    if best is None:
        return False

    paras = [p for p in best.text_frame.paragraphs if p.runs]
    while len(paras) < len(lines):
        new_p = copy.deepcopy(paras[-1]._p)
        paras[-1]._p.getparent().append(new_p)
        paras = [p for p in best.text_frame.paragraphs if p.runs]

    for i, para in enumerate(paras):
        text = lines[i] if i < len(lines) else ""
        para.runs[0].text = text
        for extra in para.runs[1:]:
            extra.text = ""
    return True


def move_slide(prs, from_idx, to_idx):
    id_list = prs.slides._sldIdLst
    entries = list(id_list)
    el = entries[from_idx]
    id_list.remove(el)
    id_list.insert(to_idx, el)


def main() -> int:
    if not PATH.exists():
        print(f"FAIL: {PATH} not found")
        return 1

    prs = Presentation(str(PATH))
    original_count = len(prs.slides)
    print(f"loaded {original_count} slides")

    hits = replace_text_everywhere(prs)
    print(f"\ntext corrections applied: {sum(hits.values())}")
    for old, n in sorted(hits.items(), key=lambda x: -x[1]):
        print(f"  {n}x  {old[:72]}")
    missed = [old for old, rep in REPLACEMENTS
              if old not in hits and old != rep]
    if missed:
        print("\nnot found (check wording):")
        for m in missed:
            print(f"  - {m[:72]}")

    made = []
    for title, data in NEW_TABLE_SLIDES:
        s = clone_slide(prs, TABLE_TEMPLATE)
        set_title(s, title)
        tbl_shape = next((sh for sh in s.shapes if sh.has_table), None)
        if tbl_shape is None:
            print(f"FAIL: cloned slide for '{title}' has no table")
            return 1
        set_table(tbl_shape, data)
        made.append(title)

    for title, lines in NEW_BULLET_SLIDES:
        s = clone_slide(prs, BULLET_TEMPLATE)
        set_title(s, title)
        if not set_bullets(s, lines):
            print(f"FAIL: cloned slide for '{title}' has no text box")
            return 1
        made.append(title)

    # New slides were appended; move them into position after INSERT_AFTER.
    for offset, _ in enumerate(made):
        move_slide(prs, original_count + offset, INSERT_AFTER + 1 + offset)

    prs.save(str(PATH))
    print(f"\nadded {len(made)} slides after slide {INSERT_AFTER + 1}:")
    for t in made:
        print(f"  {t}")
    print(f"\nsaved: {PATH}  ({original_count} -> {len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
