"""Dump the structure and text of VF_Logistics_Presentation.pptx.

Read-only. Prints slide index, layout name, and every shape with its type, position
and text so the deck can be edited precisely rather than guessed at.
"""

from pptx import Presentation
from pptx.util import Emu

PATH = (r"C:\Users\phuochoa\Mendix\VF_Logistics_Portal-main_2"
        r"\snowflake-backend\VF_Logistics_Presentation.pptx")


def main():
    prs = Presentation(PATH)
    print(f"slide size: {prs.slide_width} x {prs.slide_height} EMU "
          f"({Emu(prs.slide_width).inches:.2f} x {Emu(prs.slide_height).inches:.2f} in)")
    print(f"slides: {len(prs.slides)}")
    print("=" * 78)

    for i, slide in enumerate(prs.slides, start=1):
        print(f"\n--- SLIDE {i}  (layout: {slide.slide_layout.name}) ---")
        for shape in slide.shapes:
            kind = shape.shape_type
            name = shape.name
            ph = ""
            if shape.is_placeholder:
                ph = f" placeholder[idx={shape.placeholder_format.idx}," \
                     f"type={shape.placeholder_format.type}]"
            print(f"  <{name}> type={kind}{ph}")
            if shape.has_text_frame:
                for p_i, para in enumerate(shape.text_frame.paragraphs):
                    txt = "".join(r.text for r in para.runs)
                    if txt.strip():
                        sizes = {r.font.size.pt for r in para.runs
                                 if r.font.size is not None}
                        bolds = {r.font.bold for r in para.runs}
                        print(f"      p{p_i} lvl{para.level} "
                              f"size={sorted(sizes) if sizes else '-'} "
                              f"bold={bolds}: {txt}")
            if shape.has_table:
                tbl = shape.table
                print(f"      TABLE {len(tbl.rows)}x{len(tbl.columns)}")
                for r_i, row in enumerate(tbl.rows):
                    cells = [c.text.replace("\n", " / ") for c in row.cells]
                    print(f"        r{r_i}: {cells}")


if __name__ == "__main__":
    main()
